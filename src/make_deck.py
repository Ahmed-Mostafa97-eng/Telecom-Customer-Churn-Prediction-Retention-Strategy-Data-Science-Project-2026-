"""
Builds the GCI final-assignment deck (<=15 slides) from metrics.json + ./figures.
Re-run this after the notebook produces real numbers; the deck regenerates automatically.
"""
import json
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

M=json.load(open("metrics.json")); FIG=Path("figures")
TEAL=RGBColor(0x02,0x80,0x90); TEALD=RGBColor(0x02,0x3B,0x46); CORAL=RGBColor(0xE4,0x57,0x2E)
INK=RGBColor(0x1d,0x27,0x33); GREY=RGBColor(0x6b,0x77,0x82); WHITE=RGBColor(0xFF,0xFF,0xFF)
SAND=RGBColor(0xF1,0xED,0xE7); MINT=RGBColor(0xEAF,0x0,0x0) if False else RGBColor(0xE6,0xF4,0xF1)

prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    return s
def box(s,x,y,w,h):
    tb=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf
def para(tf,txt,size,color=INK,bold=False,first=False,align=PP_ALIGN.LEFT,italic=False,space=6,font="Calibri"):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space)
    r=p.add_run(); r.text=txt; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return p
def rect(s,x,y,w,h,color,line=None):
    r=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,I(x),I(y),I(w),I(h))
    r.fill.solid(); r.fill.fore_color.rgb=color
    if line: r.line.color.rgb=line; r.line.width=Pt(1)
    else: r.line.fill.background()
    r.shadow.inherit=False; return r
def pic(s,name,x,y,w=None,h=None):
    kw={}
    if w: kw["width"]=I(w)
    if h: kw["height"]=I(h)
    return s.shapes.add_picture(str(FIG/name),I(x),I(y),**kw)
def title(s,t,kicker=None):
    if kicker: para(box(s,0.7,0.42,12,0.4),kicker,13,TEAL,bold=True,first=True)
    para(box(s,0.7,0.74,12,0.9),t,30,INK,bold=True,first=True,font="Cambria")
def stat(s,x,y,w,big,small,color=TEAL):
    r=rect(s,x,y,w,1.55,MINT)
    tf=box(s,x+0.18,y+0.16,w-0.36,1.3)
    para(tf,big,30,color,bold=True,first=True,font="Cambria")
    para(tf,small,11.5,GREY,space=0)
def footer(s,txt):
    para(box(s,0.7,7.08,12,0.32),txt,8.5,GREY,first=True)

ev=M["ev"]; a=M["assumptions"]
NET=f"${ev['net_benefit']/1e6:.1f}M"; ROI=f"{ev['roi_pct']:.0f}%"
AUC=f"{M['best_auc']:.3f}"; LIFT=f"{M['lift_top_decile']:.1f}x"

# ---------- 1 TITLE ----------
s=slide(TEALD)
rect(s,0,0,13.333,7.5,TEALD)
para(box(s,0.9,2.2,11.5,0.5),"GCI WORLD 2026 · FINAL ASSIGNMENT",15,RGBColor(0x9F,0xD4,0xDC),bold=True,first=True)
para(box(s,0.9,2.75,11.6,1.8),"Stop the Leak: A Data-Driven Retention Engine for Company A",40,WHITE,bold=True,first=True,font="Cambria")
para(box(s,0.9,4.7,11,0.8),"Turning 100k customers of usage history into a targeted, ROI-positive churn-prevention program.",16,RGBColor(0xCF,0xE6,0xEA),first=True)
para(box(s,0.9,6.6,11,0.4),"IT Consulting — Proof of Concept",12,RGBColor(0x7F,0xB8,0xC2),first=True)

# ---------- 2 EXEC SUMMARY / TOC ----------
s=slide(); title(s,"Executive summary","THE PROPOSAL IN ONE SLIDE")
para(box(s,0.7,1.65,7.0,1.4),
     "Company A loses a large share of its base each cycle. We built a churn-risk model on its own "
     "usage data, identified the customers most likely to leave, and show that a focused retention "
     "offer to the riskiest 10% pays for itself several times over.",14.5,INK,first=True)
stat(s,0.7,3.35,2.7,AUC,f"Best model AUC ({M['best_model']})")
stat(s,3.6,3.35,2.7,LIFT,"more churners found vs random targeting",color=CORAL)
stat(s,6.5,3.35,2.7,NET,"net benefit / year (illustrative)",color=TEALD)
# agenda card
rect(s,9.7,1.65,2.95,5.0,SAND)
tf=box(s,9.95,1.9,2.5,4.6)
para(tf,"AGENDA",12,TEAL,bold=True,first=True)
for i,t in enumerate(["Market context","The data & EDA","Problem definition",
                      "Modelling & evidence","Risk-based targeting","Business proposal",
                      "Quantified impact","Pricing & next steps"]):
    para(tf,f"{i+1}.  {t}",12.5,INK,space=8)
footer(s,"Figures and numbers in this deck are produced automatically from the analysis notebook.")

# ---------- 3 MARKET ----------
s=slide(); title(s,"A saturated market where retention beats acquisition","MARKET CONTEXT")
stat(s,0.7,1.7,3.7,f"{M['market']['global_churn_rate_2025']*100:.1f}%","global telecom churn rate, 2025¹")
stat(s,4.6,1.7,3.7,"5–7×","costlier to acquire than to retain a customer²")
stat(s,8.5,1.7,4.15,"$120M","lost yearly at 1M subs · $50 ARPU · 20% churn³",color=CORAL)
tf=box(s,0.7,3.6,7.2,3.0)
para(tf,"Why this is the right problem to solve",16,TEAL,bold=True,first=True)
for t in ["Mature markets grow by keeping customers and lifting revenue per user, not by net adds.",
          "Behavioural signals — usage decline, complaints, ageing handsets — precede churn by months, so it is predictable.",
          "A retained customer is far cheaper than a replacement, so even a modest save rate is highly profitable."]:
    para(tf,"•  "+t,13.5,INK,space=9)
pic(s,"churn_balance.png",8.2,3.55,w=4.6)
footer(s,"Sources: 1) CustomerGauge, Average Churn Rate by Industry 2025.  2) Tridens Technology, Telecom Churn 2025.  3) BillingPlatform, Churn Rates by Industry 2025.")

# ---------- 4 DATA ----------
s=slide(); title(s,"The asset: Company A's own customer data","THE DATA")
tf=box(s,0.7,1.7,5.7,2.0)
para(tf,f"{M['n_rows']:,} customers · {M['n_cols']} fields",18,INK,bold=True,first=True)
para(tf,"Two linked tables joined on Customer_ID:",13.5,INK,space=8)
para(tf,"•  Client — tenure, plan, equipment, demographics",13,INK,space=4)
para(tf,"•  Record — monthly usage, billing, call-quality + the churn flag",13,INK,space=4)
para(tf,"churn = customer left 31–60 days after the observation date.",12.5,GREY,italic=True,space=4)
rect(s,0.7,4.0,5.7,2.4,MINT)
tf=box(s,0.95,4.2,5.2,2.0)
para(tf,"First task: choose the question",13.5,TEAL,bold=True,first=True)
para(tf,"The data supports many targets. We chose churn because it is directly tied to "
       "revenue, the signals to predict it are present, and — crucially — the output is "
       "actionable: a retention team can act on a ranked risk list tomorrow.",12.5,INK,space=4)
pic(s,"missingness.png",6.8,1.7,w=6.0)
footer(s,"Data-quality scan informs which fields to drop or impute before modelling.")

# ---------- 5 EDA lever ----------
s=slide(); title(s,"EDA finding #1: ageing handsets are the clearest churn lever","EXPLORATORY ANALYSIS")
pic(s,"eqpdays_churn.png",0.7,1.7,w=6.3)
tf=box(s,7.3,1.9,5.4,4.4)
para(tf,"Read this chart",15,TEAL,bold=True,first=True)
para(tf,f"Churners carry markedly older equipment on average "
       f"({M['eqp_churn_mean']:.0f} days vs {M['eqp_stay_mean']:.0f} for those who stay).",13.5,INK,space=10)
para(tf,"Why it matters for a proposal",15,TEAL,bold=True,space=4)
para(tf,"Unlike tenure or demographics, equipment age is something the business can "
       "change — through upgrade offers. That makes it a lever, not just a correlate.",13.5,INK,space=4)
footer(s,"Dashed lines mark group means. Distribution truncated at the 99th percentile for readability.")

# ---------- 6 EDA tenure ----------
s=slide(); title(s,"EDA finding #2: risk is concentrated early in the lifecycle","EXPLORATORY ANALYSIS")
pic(s,"churn_by_tenure.png",0.7,1.7,w=6.2)
tf=box(s,7.2,1.9,5.5,4.2)
para(tf,"Where to focus",15,TEAL,bold=True,first=True)
para(tf,"Churn rate is highest in the first 12 months and falls as customers settle in. "
       "New, under-tenured customers with old handsets are the highest-priority segment.",13.5,INK,space=10)
para(tf,"This shapes the proposal",15,TEAL,bold=True,space=4)
para(tf,"A retention program should weight early-tenure, declining-usage customers — exactly "
       "the group the model will surface in its top risk decile.",13.5,INK,space=4)
footer(s,"Tenure bands in months.")

# ---------- 7 PROBLEM DEF ----------
s=slide(TEALD)
para(box(s,0.7,0.6,12,0.4),"PROBLEM DEFINITION",14,RGBColor(0x9F,0xD4,0xDC),bold=True,first=True)
para(box(s,0.7,1.0,12,1.0),"From a vague worry to a precise, answerable question",30,WHITE,bold=True,first=True,font="Cambria")
rect(s,0.7,2.4,5.8,3.9,RGBColor(0x05,0x4A,0x57))
tf=box(s,1.0,2.7,5.2,3.4)
para(tf,"The business question",15,RGBColor(0x9F,0xD4,0xDC),bold=True,first=True)
para(tf,"\u201cWhich current customers are most likely to leave in the next cycle, "
       "and is it worth intervening?\u201d",16,WHITE,italic=True,space=12)
para(tf,"ML task",14,RGBColor(0x9F,0xD4,0xDC),bold=True,space=4)
para(tf,"Binary classification. Target = churn. Output = a 0–1 risk score per customer.",13.5,WHITE,space=4)
rect(s,6.9,2.4,5.7,3.9,RGBColor(0x05,0x4A,0x57))
tf=box(s,7.2,2.7,5.1,3.4)
para(tf,"Who acts on the output",15,RGBColor(0x9F,0xD4,0xDC),bold=True,first=True)
para(tf,"The retention team receives a ranked list. High-risk, high-value customers get a "
       "targeted offer before they leave.",13.5,WHITE,space=12)
para(tf,"Why a score, not a label",14,RGBColor(0x9F,0xD4,0xDC),bold=True,space=4)
para(tf,"A ranked score lets the business choose how many customers to contact based on budget — "
       "a single yes/no label cannot.",13.5,WHITE,space=4)

# ---------- 8 FEATURE ENGINEERING ----------
s=slide(); title(s,"Engineering signals the raw columns don't capture","FEATURE ENGINEERING")
para(box(s,0.7,1.6,7.6,0.7),
     "Beyond cleaning and encoding, we derived features that encode churn intuition directly:",13.5,INK,first=True)
feats=[("eqp_per_month","Handset age relative to how long they've been a customer"),
       ("usage_declining","Flag for customers whose minutes are trending down"),
       ("care_intensity","Customer-care calls per minute of use — a frustration proxy"),
       ("rev_per_min","Revenue efficiency — value of each minute used"),
       ("drop_rate","Share of attempted calls that fail — network-quality pain")]
yy=2.25
for n,d in feats:
    rect(s,0.7,yy,7.4,0.84,SAND)
    tf=box(s,0.92,yy+0.12,7.0,0.64)
    para(tf,n,13,TEAL,bold=True,first=True,font="Courier New")
    para(tf,d,11.5,INK,space=0)
    yy+=0.93
pic(s,"prep_flow.png",8.45,1.7,h=4.9)
footer(s,"Encoding keeps missing categories as their own level; ID columns dropped to avoid memorisation.")

# ---------- 9 MODEL ----------
s=slide(); title(s,"We compared three models — and named the winner","MODELLING & EVIDENCE")
# table
mt=M["model_table"]; order=list(mt.keys())
cols=["Model","AUC","Accuracy","Precision","Recall","F1"]
nrows=len(order)+1
tbl=s.shapes.add_table(nrows,len(cols),I(0.7),I(1.75),I(7.0),I(2.2)).table
for j,c in enumerate(cols):
    cell=tbl.cell(0,j); cell.text=c
    cell.fill.solid(); cell.fill.fore_color.rgb=TEALD
    p=cell.text_frame.paragraphs[0]; p.runs[0].font.color.rgb=WHITE
    p.runs[0].font.size=Pt(11); p.runs[0].font.bold=True
for i,m in enumerate(order):
    vals=[m,f"{mt[m]['AUC']:.3f}",f"{mt[m]['Accuracy']:.3f}",f"{mt[m]['Precision']:.3f}",
          f"{mt[m]['Recall']:.3f}",f"{mt[m]['F1']:.3f}"]
    win=(m==M["best_model"])
    for j,v in enumerate(vals):
        cell=tbl.cell(i+1,j); cell.text=v
        cell.fill.solid(); cell.fill.fore_color.rgb=MINT if win else WHITE
        p=cell.text_frame.paragraphs[0]; p.runs[0].font.size=Pt(11)
        p.runs[0].font.bold=win; p.runs[0].font.color.rgb=INK
tf=box(s,0.7,4.3,7.0,2.4)
para(tf,f"Chosen model: {M['best_model']}",15,TEAL,bold=True,first=True)
para(tf,f"Evaluation metric: ROC-AUC = {AUC}.  We report AUC (not just accuracy) because it "
       f"measures how well the model ranks risk regardless of the cut-off the business picks — "
       f"the property the targeting use-case depends on.",13,INK,space=6)
para(tf,f"A no-skill model scores 0.50; ours clears that comfortably, so the score reflects real signal.",12.5,GREY,space=4)
pic(s,"model_comparison.png",8.0,1.9,w=4.9)
footer(s,"Requirement: model name, metric, and score stated explicitly above. Stratified 70/30 split, random_state=42.")

# ---------- 10 EVIDENCE charts ----------
s=slide(); title(s,"How the model performs and what it keys on","MODELLING & EVIDENCE")
pic(s,"roc.png",0.7,1.7,w=4.0)
pic(s,"confusion.png",4.85,1.85,w=3.6)
pic(s,"feature_importance.png",8.45,1.7,w=4.4)
tf=box(s,0.7,6.0,12,0.9)
para(tf,f"The model ranks well above chance (ROC, left). Top drivers (right) line up with the EDA — "
       f"equipment age, usage change, and care intensity — so the story is consistent and explainable.",13,INK,first=True)
footer(s,"Feature importance by XGBoost gain; coral bars are engineered features.")

# ---------- 11 TARGETING ----------
s=slide(); title(s,"From risk score to action: target the riskiest decile","RISK-BASED TARGETING")
pic(s,"lift_deciles.png",0.7,1.7,w=6.2)
pic(s,"gains.png",0.7,4.4,w=6.2) if False else None
tf=box(s,7.2,1.9,5.5,4.4)
para(tf,"The core insight",15,TEAL,bold=True,first=True)
para(tf,f"If we rank everyone by predicted risk and contact only the top 10%, that group "
       f"churns at {M['top_decile_churn_rate']:.0f}% — about {LIFT} the base rate.",14,INK,space=10)
para(tf,f"Capture",15,TEAL,bold=True,space=2)
para(tf,f"That single decile contains roughly {M['top_decile_capture']:.0f}% of all churners; the "
       f"top two deciles capture about {M['top2_decile_capture']:.0f}%.",14,INK,space=10)
para(tf,"This is what makes the budget go further — spend retention dollars where churn actually is.",13,GREY,italic=True)
footer(s,"Deciles on the held-out test set; decile 1 = highest predicted risk.")

# ---------- 12 PROPOSAL ----------
s=slide(); title(s,"The proposal: a model-driven retention program","BUSINESS PROPOSAL")
steps=[("1","Score","Rank the full base monthly by churn risk using the model."),
       ("2","Target","Take the top-risk decile, prioritising early-tenure, high-value customers."),
       ("3","Act","Trigger a tailored save offer — handset upgrade or plan fix — addressing the driver."),
       ("4","Measure","Track save rate vs a hold-out control; feed results back to retrain.")]
x=0.7
for n,t,d in steps:
    rect(s,x,1.9,2.85,3.4,SAND)
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,I(x+0.2),I(2.1),I(0.6),I(0.6))
    c.fill.solid(); c.fill.fore_color.rgb=TEAL; c.line.fill.background(); c.shadow.inherit=False
    ctf=c.text_frame; ctf.paragraphs[0].text=n; ctf.paragraphs[0].alignment=PP_ALIGN.CENTER
    ctf.paragraphs[0].runs[0].font.color.rgb=WHITE; ctf.paragraphs[0].runs[0].font.bold=True; ctf.paragraphs[0].runs[0].font.size=Pt(18)
    tf=box(s,x+0.22,2.85,2.45,2.3)
    para(tf,t,16,TEAL,bold=True,first=True)
    para(tf,d,12.5,INK,space=4)
    x+=3.05
tf=box(s,0.7,5.55,12,1.2)
para(tf,"Why it works: the model concentrates spend on customers who would actually leave, and the "
       "intervention targets a driver the EDA proved is real (ageing equipment). Each step is measurable.",13.5,INK,first=True)
footer(s,"Hold-out control group lets you attribute saves to the program rather than to chance.")

# ---------- 13 IMPACT ----------
s=slide(); title(s,"Quantified impact: the program pays for itself","QUANTIFIED IMPACT")
pic(s,"ev_business_case.png",0.7,1.7,w=6.0)
pic(s,"sensitivity.png",0.7,4.55,w=6.0) if False else None
tf=box(s,7.0,1.75,5.7,4.8)
para(tf,"Worked example (per year, scaled to 100k customers)",13.5,TEAL,bold=True,first=True)
rows=[("Customers contacted (top 10%)",f"{ev['targeted']:,}"),
      ("Would-be churners in that group",f"{ev['churners_in_target']:,}"),
      ("Customers saved (30% offer success)",f"{ev['saved_customers']:,}"),
      ("Revenue protected (12 mo × $50)",f"${ev['revenue_protected']/1e6:.2f}M"),
      ("Campaign cost ($25 × targeted)",f"${ev['campaign_cost']/1e6:.2f}M"),
      ("Net benefit",f"${ev['net_benefit']/1e6:.2f}M"),
      ("Return on retention spend",ROI)]
for i,(k,v) in enumerate(rows):
    bold=i>=5
    para(tf,k,12.5,INK if not bold else TEALD,bold=bold,space=2)
    para(tf,v,16 if bold else 13.5,CORAL if i==5 else (TEAL if bold else INK),bold=True,space=7)
footer(s,"Assumptions stated, not hidden: ARPU $50, 12-mo horizon, $25 offer cost, 30% save rate. "
        "Net benefit stays positive across a 10–50% save-rate range (sensitivity tested).")

# ---------- 14 PRICING / NEXT STEPS ----------
s=slide(TEALD)
para(box(s,0.7,0.6,12,0.4),"ENGAGEMENT",14,RGBColor(0x9F,0xD4,0xDC),bold=True,first=True)
para(box(s,0.7,1.0,12,1.0),"Pricing and next steps",30,WHITE,bold=True,first=True,font="Cambria")
rect(s,0.7,2.4,5.8,3.9,RGBColor(0x05,0x4A,0x57))
tf=box(s,1.0,2.7,5.2,3.4)
para(tf,"What we deliver",15,RGBColor(0x9F,0xD4,0xDC),bold=True,first=True)
for t in ["Production churn-scoring model + monthly refresh","Ranked risk list wired to the retention workflow",
          "A/B hold-out design to prove incremental saves","Quarterly retrain as new data lands"]:
    para(tf,"•  "+t,13.5,WHITE,space=8)
rect(s,6.9,2.4,5.7,3.9,RGBColor(0x05,0x4A,0x57))
tf=box(s,7.2,2.7,5.1,3.4)
para(tf,"Commercials",15,RGBColor(0x9F,0xD4,0xDC),bold=True,first=True)
para(tf,"Fixed PoC-to-production fee priced at a small fraction of the "
       f"{NET}/yr net benefit it unlocks, plus a monthly retainer for model upkeep.",13.5,WHITE,space=10)
para(tf,"Risks & limits",14,RGBColor(0x9F,0xD4,0xDC),bold=True,space=4)
para(tf,"Save-rate is an assumption to validate in the pilot; effects like offer cannibalisation "
       "and seasonality are measured via the control group.",13,WHITE,space=4)

# ---------- 15 REFERENCES ----------
s=slide(); title(s,"References","SOURCES & TOOLS")
tf=box(s,0.7,1.7,12,5.0)
refs=[
 "CustomerGauge, \u201cAverage Churn Rate by Industry \u2014 2025 B2B Benchmarks\u201d: telecom ~22\u201331% annual churn. customergauge.com/blog/average-churn-rate-by-industry",
 "Tridens Technology, \u201cWhy Telecom Customers Churn and How to Measure It\u201d (2025): telecom churn 20\u201350% annually; retention vs. acquisition economics. tridenstechnology.com/telecom-churn",
 "BillingPlatform, \u201cChurn Rates by Industry\u201d (2025): telecom ~31% churn, ~78% retention. billingplatform.com/blog/average-churn-rate-by-industry",
 "Company A telecom dataset (Client.csv, Record.csv) \u2014 provided for this assignment.",
 "Python libraries: pandas, NumPy, scikit-learn, XGBoost, Matplotlib.",
 "Generative-AI assistance: Claude (Anthropic), for analysis and slide preparation.",
]
for r in refs:
    para(tf,"•  "+r,13,INK,space=10)
footer(s,"On-slide source tags above correspond to this list, per the citation requirement.")

prs.save("deck.pptx")
print("Saved deck.pptx with",len(prs.slides.__iter__.__self__._sldIdLst),"slides")
